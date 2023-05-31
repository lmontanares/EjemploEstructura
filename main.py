import collections
import collections.abc
import json
import os
import time
from pathlib import Path
from zipfile import ZipFile as zipfile
import pptx as pptx
import requests
from dotenv import load_dotenv
from jinja2 import Environment, FileSystemLoader
from loguru import logger
from netmiko import ConnectHandler
from playwright.sync_api import sync_playwright
from pptx import Presentation
from models import MensajeRCS, MensajeRCS_v2
from rcs_bot import create_rcs_campaign, login, upload_images
from messages import OpenUrlAction, DialAction, StandaloneCard
import unicodedata
import re
import unidecode

load_dotenv()


def unzip_files():
    files = Path(__file__).parent.glob("data/*.zip")
    if files:
        for file in files:
            zip_ref = zipfile(file, "r")
            zip_ref.extractall(Path(__file__).parent / "data")
            zip_ref.close()


def read_pptx() -> list:
    """
    Lee todos los archivos PowerPoint (.pptx) de la carpeta "data" y devuelve una lista de tablas como cadenas de texto.

    Returns:
    - Lista de cadenas de texto, cada una representando una tabla de un archivo PowerPoint.
    """
    # Obtener todos los archivos .pptx en la carpeta "files".
    files = Path(__file__).parent.glob("data/*.pptx")

    # Inicializar una lista vacía para almacenar las tablas del archivo.
    tables = []
    # Iterar sobre cada archivo .pptx.
    for file in files:
        # Crear un objeto `Presentation` a partir del archivo.
        pres = Presentation(file)
        # Iterar sobre las diapositivas del archivo.
        for slide in pres.slides:
            # Iterar sobre los objetos de forma de la diapositiva.
            for shape in iter(slide.shapes):
                # Si el objeto de forma es una tabla, agregarla a la lista de tablas.
                if shape.has_table:
                    table = shape.table
                    # Agregar la tabla a la lista, omitiendo las celdas impares.
                    tables.append([cell.text.strip() for i, cell in enumerate(table.iter_cells()) if i % 2 != 0])

        # Convertir cada tabla en una cadena de texto.
        for i, table in enumerate(tables):
            tables[i] = "|" + "|".join(table)  # type: ignore
            # tables[i] = str(["|" + item for item in table])
    return tables


def remove_special_characters(string: str):
    return re.sub(r"[^a-zA-Z0-9_]", "", string)


def normalize_data(data: list[str]) -> list[str]:
    normalized_data = []
    for item in data:
        cleaned_item = re.sub(r"[\u200B\u200C\u200D\u2060\uFEFF]", "", item)
        normalized_item = [unicodedata.normalize("NFKD", x).replace("\n", " ").strip() for x in cleaned_item.split("|")]
        normalized_data.append(normalized_item)
    return normalized_data


def define_post_back_data(text: str, type: str):
    if text:
        match type:
            case "OpenUrlAction":
                postback_data = f"URL_{text}"
            case "DialAction":
                postback_data = f"DIAL_{text}"
            case _:
                pass

        postback_data = unidecode.unidecode(text.title())
        postback_data = remove_special_characters(postback_data)
        postback_data = f"{postback_data}".replace(" ", "")

    redirection_id = input(
        f'Enter campaign id for redirection of Button "{text}" , or press enter to default "{postback_data}": '
    )

    return redirection_id if redirection_id else postback_data


def pptx_data_to_stand_alone_dict(data: list[list[str]]) -> list[dict[str, str | list[str]]]:
    msg_data = []
    for row in data:
        suggestions = [row[i : i + 2] for i in range(5, len(row) - 1, 2)]
        dict_data = {
            "id": row[0],
            "name": row[1],
            "filename": row[2].replace(" ", "_").strip(),
            "title": row[3],
            "msg": row[4],
            "fallback_msg": row[-1],
            "suggestions": suggestions,
        }
        msg_data.append(dict_data)
    return msg_data


def stand_alone_dict_to_class(msg_data: list[dict[str, str | list[str]]]) -> list[MensajeRCS_v2]:
    stand_alone_list = []
    for msg in msg_data:
        if not msg["id"]:
            input_ = input(f"Introduce el ID del mensaje {msg['name']}: ")
            msg["id"] = input_.strip()
        suggestions = []
        for suggestion in msg["suggestions"]:
            if suggestion[0] != "":
                # Cliente envia este numero con un asterisco (*) al inicio, por eso se hace el slice
                if suggestion[-1][1:].isnumeric():
                    postback_data = define_post_back_data(suggestion[0], "DialAction")
                    suggestions.append(
                        DialAction(
                            text=suggestion[0],
                            phone_number=suggestion[1][1:],
                            postback_data=postback_data,
                        )
                    )
                else:
                    postback_data = define_post_back_data(suggestion[0], "OpenUrlAction")
                    suggestions.append(
                        OpenUrlAction(text=suggestion[0], url=suggestion[1], postback_data=postback_data)
                    )
        stand_alone_list.append(
            MensajeRCS_v2(
                id=msg["id"],
                name=msg["name"],
                # unidecode.unidecode(str) remove acutes, accents and other special characters
                fallback_msg=unidecode.unidecode(msg["fallback_msg"]),
                url=None,
                rcs_msg=StandaloneCard(
                    card_orientation="VERTICAL",
                    media_height="MEDIUM",
                    file_resource_id=f'{os.environ.get("STATIC_RCS_URL")}{msg["filename"]}',
                    title=msg["title"],
                    description=msg["msg"],
                    suggestions=suggestions,
                ),
            )
        )
    return stand_alone_list


def get_msg_data_stand_alone(pptx_data: list[str]):
    data = []
    for x in pptx_data:
        x_split = x.split("|")
        dict_data = dict(
            zip(
                [
                    "id",
                    "name",
                    "filename",
                    "title_rcs_msg",
                    "rcs_msg",
                    "button1_text",
                    "button1_url",
                    "button2_text",
                    "button2_url",
                    "fallback_msg",
                ],
                x_split,
            )
        )
        data.append(dict_data)
    return data


def create_stand_alone(data: list) -> list[MensajeRCS]:
    msg_list = []
    for x in data:
        msg_list.append(MensajeRCS(**x))
    return msg_list


def render_jinja(msg_data: dict) -> str:
    """
    Renderiza una plantilla Jinja a partir de un diccionario de datos.

    Parameters:
    - msg_data: El diccionario de datos con el que se rellenará la plantilla.

    Returns:
    El resultado del renderizado de la plantilla con los datos proporcionados.
    """
    # Crear un objeto `Environment` con el cargador de plantillas de archivo.
    env = Environment(loader=FileSystemLoader(Path(__file__).parent / "templates"))

    # Obtener la plantilla con el nombre "json_template.j2".
    template = env.get_template("json_template.j2")

    # Renderizar la plantilla con los datos proporcionados.
    output = template.render(msg_data)

    # Devolver el resultado del renderizado.
    return output


def delete_campaign():
    # Delete campaign code goes here
    print("Deleting campaign...")


def send_rcs_message(id: str, msisdn: str = "52000000001"):
    """
    Envía un mensaje RCS a un número de teléfono específico.

    Parameters:
    - id: id del mensaje RCS a enviar
    - msisdn: número de teléfono al que se enviará el mensaje (opcional, default es 52000000001)

    Returns:
    - Resultado de la solicitud de envío de mensaje RCS.
    """
    requests.post(
        url=os.environ.get("SEND_RCS_URL"),
        json={
            "account_id": os.environ.get("ACCOUNT_ID"),
            "account_key": os.environ.get("ACCOUNT_KEY"),
            "msisdn": msisdn,
            "transaction_id": "pruebas_cliente",
            "rcs_message_id": id,
            "campaign_id": "pruebas_cliente",
            "eta": "now",
            "validity": "720",
            "priority": "0",
            "ticket_id": "666",
            "ctx": {},
        },
    )
    logger.info(f"Sending {id} to {msisdn}.")


def get_urls(msg_list: list[MensajeRCS_v2], msisdn: str = "52000000001"):
    """
    Obtiene las URLs de los mensajes RCS enviados a un número de teléfono específico.

    Parameters:
    - msg_list: lista de objetos MensajeRCS para los cuales se buscarán las URLs.
    - msisdn: número de teléfono para el cual se buscarán las URLs (opcional, default es 52000000001)

    Returns:
    - Lista de URLs de los mensajes RCS enviados al número de teléfono especificado.
    """
    # Obtener valores de las variables de entorno del sistema.
    ips = [
        # os.environ.get("SERVER1"),
        # os.environ.get("SERVER2"),
        # os.environ.get("SERVER3"),
        os.environ.get("SERVER4"),
        os.environ.get("SERVER5"),
        os.environ.get("SERVER6"),
        # os.environ.get("SERVER7"),
    ]
    for i, ip in enumerate(ips):
        device = {
            "device_type": "linux",
            "ip": ip,
            "username": "root",
            # "password": f"{os.environ.get('PASSWORD') if i < 3 else os.environ.get('PASSWORD2')}",
            "password": os.environ.get("PASSWORD_2"),
            # "port": f"{22 if i < 3 else 2222}",
            "port": 2222,
        }

        with ConnectHandler(**device) as connection:
            for msg in msg_list:
                if msg.id:
                    cmd = (
                        rf"grep mvads /var/log/antica-vas/webapp.log  | grep Sending | grep {msg.id} | grep {msisdn}| awk -F '"
                        " "
                        "' '{print $NF}'"
                    )
                    output = connection.send_command(cmd)
                    # urls = [x if "https" in x else "None" for x in output.split()]
                    urls = [x for x in output.split()]
                    if urls:
                        msg.url = urls

    return msg_list


def main():
    unzip_files()
    data = read_pptx()
    data = normalize_data(data)
    data = pptx_data_to_stand_alone_dict(data)
    msg_list = stand_alone_dict_to_class(data)

    # Si se desea que el navegador se abra, cambiar headless a False.

    with sync_playwright() as playwright:
        page, browser, context = login(playwright, headless=False)
        create_rcs_campaign(page, msg_list)
        upload_images(page)
        context.close()
        browser.close()

    for msg in msg_list:
        if msg.id:
            send_rcs_message(msg.id, msisdn="52000000001")

    msg_list = get_urls(msg_list, msisdn="52000000001")

    for msg in msg_list:
        if msg.url:
            logger.info(f"Campaign={msg.id}|ULRs={msg.url}")

    print("Message List for EXCEL")
    for msg in msg_list:
        print(f"{msg.id}|{msg.name}")


if __name__ == "__main__":
    main()
